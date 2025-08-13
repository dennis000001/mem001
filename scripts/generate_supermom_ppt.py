from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def add_bullets(text_frame, lines):
    """
    Add bullet points to a text frame.
    `lines` is a list of dicts: {"text": str, "level": int}
    """
    # Ensure first paragraph uses the placeholder's first paragraph
    first = True
    for item in lines:
        if first:
            p = text_frame.paragraphs[0]
            p.text = item["text"]
            p.level = item.get("level", 0)
            first = False
        else:
            p = text_frame.add_paragraph()
            p.text = item["text"]
            p.level = item.get("level", 0)
        p.space_after = Pt(6)


def make_presentation(output_path: str):
    prs = Presentation()

    # Common layout: Title and Content
    TITLE_AND_CONTENT = 1

    # Slide 1: 愿景与北极星
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    slide.shapes.title.text = "01 愿景与北极星"
    tf = slide.shapes.placeholders[1].text_frame
    bullets = [
        {"text": "愿景：用 AI 服务一百万位妈妈，让每次选择更省时、省钱、更安心、更贴心", "level": 0},
        {"text": "北极星：Mom Value Index（MVI）= 省时 + 省钱 + 安心 + 贴心（加权）", "level": 0},
        {"text": "三年目标：服务≥100万；时间成本↓≥30%；人均月节省≥5%；可信命中≥85%；福利触达≥60%", "level": 0},
    ]
    add_bullets(tf, bullets)

    # Slide 2: 市场与定位
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    slide.shapes.title.text = "02 市场与定位"
    tf = slide.shapes.placeholders[1].text_frame
    bullets = [
        {"text": "市场机会：东南亚6国，母婴消费升级，线上+O2O并进", "level": 0},
        {"text": "平台定位：母婴数据与AI平台（零方数据×洞察×营销激活×C端体验）", "level": 0},
        {"text": "时机：隐私趋严→零/一方数据为王；多语多文化带来AI本地化红利", "level": 0},
    ]
    add_bullets(tf, bullets)

    # Slide 3: 现状与能力
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    slide.shapes.title.text = "03 现状与能力"
    tf = slide.shapes.placeholders[1].text_frame
    bullets = [
        {"text": "数据与AI：零方数据（阶段/偏好/问卷/活动）、家庭图谱、RAG可信内容、个性化排序", "level": 0},
        {"text": "标志能力：极速调研（千人≈1小时、可跨国并行）、价格/返利提醒、省钱路径规划", "level": 0},
        {"text": "触点与合作：社区/导购/试用/展会/O2O；对接 Lazada/Shopee/Watsons 等、200+ 品牌", "level": 0},
        {"text": "参考：TechNode（2023/2024）、36氪（极速调研）", "level": 0},
    ]
    add_bullets(tf, bullets)

    # Slide 4: 竞品对比
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    slide.shapes.title.text = "04 竞品对比"
    tf = slide.shapes.placeholders[1].text_frame
    bullets = [
        {"text": "theAsianparent/Parentinc：强社区与D2C；启示：做强‘可信内容+工具’留存引擎", "level": 0},
        {"text": "Shopee Moms Club：强站内个性化与价格体系；启示：打造跨平台‘省钱+决策助理’", "level": 0},
        {"text": "Mummys Market：强线下转化；启示：逛展数据化与闭环回收", "level": 0},
        {"text": "SuperMom差异：深度零方数据+快速洞察/增量测量+O2O真实世界反馈", "level": 0},
    ]
    add_bullets(tf, bullets)

    # Slide 5: 三年路线图（0–36月合并）
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    slide.shapes.title.text = "05 三年路线图（0–36月合并）"
    tf = slide.shapes.placeholders[1].text_frame
    bullets = [
        {"text": "0–9月（MVP与底座）", "level": 0},
        {"text": "阶段化首页；语义搜索+RAG；价格降价+返利；本地福利卡；社区质检；多语基线（EN/ID）", "level": 1},
        {"text": "KPI：TTF≤20s；价格提醒设置率≥15%；月节省≥1.5%", "level": 1},
        {"text": "9–18月（跨语/增量/O2O）", "level": 0},
        {"text": "逛展助理；试用/众测引擎；补货提醒；跨语语义对齐；uplift实验平台；清洁室", "level": 1},
        {"text": "KPI：月节省≥3%；到场率≥35%；可信命中≥85%；试用完成≥60%", "level": 1},
        {"text": "18–36月（平台化与生态）", "level": 0},
        {"text": "一体化AI助理；KOC/KOL协同；医疗机构非诊疗合作；LTV与预算优化；隐私增强", "level": 1},
        {"text": "KPI：活跃妈妈≥100万；月节省≥5%；福利触达≥60%；MVI持续提升", "level": 1},
    ]
    add_bullets(tf, bullets)

    # Slide 6: 风险与对策
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    slide.shapes.title.text = "06 风险与对策"
    tf = slide.shapes.placeholders[1].text_frame
    bullets = [
        {"text": "代表性偏差：配额/加权/多源采集", "level": 0},
        {"text": "医疗敏感：红线词+RAG权威来源+人工复核", "level": 0},
        {"text": "多语长尾：模型蒸馏+人审+社区纠错激励", "level": 0},
        {"text": "隐私合规：同意/撤回/最小化/保留期；清洁室/联邦学习", "level": 0},
        {"text": "归因难度：uplift与实验化框架（geo/时间分层）", "level": 0},
    ]
    add_bullets(tf, bullets)

    # Slide 7: 组织与协作
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    slide.shapes.title.text = "07 组织与协作"
    tf = slide.shapes.placeholders[1].text_frame
    bullets = [
        {"text": "三小队：Search & Trust（省时/安心）/ Deals & Commerce（省钱）/ Care & Community（贴心）", "level": 0},
        {"text": "两平台：Data & AI / Trust & Safety", "level": 0},
        {"text": "能力栈：RAG/NLP、多语语义、排序推荐、价格与优惠解析、反作弊、清洁室与实验平台", "level": 0},
        {"text": "OKR对齐：MVI、节省率、可信命中、活动到场、试用完成", "level": 0},
    ]
    add_bullets(tf, bullets)

    # Slide 8: 行动号召
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    slide.shapes.title.text = "08 行动号召"
    tf = slide.shapes.placeholders[1].text_frame
    bullets = [
        {"text": "90天交付：阶段化首页+RAG V1+价格提醒+福利卡+埋点仪表板", "level": 0},
        {"text": "资源与合作：数据/搜索/推荐/前端/合规/BD；电商/线下/医疗伙伴共创", "level": 0},
        {"text": "验收口径：指标达标+NPS访谈+品牌复购与增量显著性", "level": 0},
        {"text": "下一步：确认PoC国家/类目、里程碑排期与Owner、共创名单", "level": 0},
    ]
    add_bullets(tf, bullets)

    prs.save(output_path)


if __name__ == "__main__":
    output = "/workspace/SuperMom_AI_for_1M_Moms_CN.pptx"
    make_presentation(output)
    print(f"Generated: {output}")